# presentation_generator.py - Generate PowerPoint presentations using ChatGPT
import os
import re
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configure OpenAI API
openai.api_key = os.environ.get('OPENAI_API_KEY')

def generate_presentation(text_content, output_file, title="Generated Presentation", style="professional"):
    """Generate a PowerPoint presentation from text content."""
    # Get slide content from ChatGPT
    slides_content = get_slides_from_chatgpt(text_content, title, style)
    
    # Create the presentation
    prs = Presentation()
    
    # Apply presentation style
    apply_presentation_style(prs, style)
    
    # Add title slide
    create_title_slide(prs, title, style)
    
    # Create content slides
    for slide_data in slides_content:
        create_content_slide(prs, slide_data, style)
    
    # Save the presentation
    prs.save(output_file)
    return output_file

def get_slides_from_chatgpt(text_content, title, style):
    """Get slide content suggestions from ChatGPT."""
    # Prepare prompt for ChatGPT
    prompt = f"""
    I need to convert the following text into a well-structured PowerPoint presentation.
    The presentation title is: "{title}"
    Style/Tone: {style}
    
    Please analyze the text below and create a structured presentation outline with:
    1. A title slide
    2. An introduction/overview slide
    3. Several content slides with main points and supporting details
    4. A conclusion/summary slide
    
    For each slide, provide:
    - A clear, concise heading (maximum 7 words)
    - Bullet points for key content (2-5 points per slide, each point should be brief)
    - Any visualization suggestions (charts, graphs, images)
    
    Text to convert:
    {text_content[:6000]}  # Limit text to avoid token limits
    
    Format your response as a JSON array of slide objects with 'title', 'points' (array), and 'visual_suggestion' properties.
    """
    
    # Call ChatGPT API
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",  # or "gpt-4" for better results
            messages=[
                {"role": "system", "content": "You are a presentation expert that converts text to well-structured PowerPoint slides. Respond only with the requested JSON format."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        # Extract and parse JSON response
        content = response.choices[0].message.content
        
        # Find JSON content (if embedded in explanatory text)
        json_match = re.search(r'\[\s*{.*}\s*\]', content, re.DOTALL)
        if json_match:
            import json
            try:
                slides_content = json.loads(json_match.group(0))
                return slides_content
            except json.JSONDecodeError:
                pass
        
        # Fallback to simple parsing if JSON extraction fails
        slides = []
        lines = content.strip().split('\n')
        current_slide = None
        
        for line in lines:
            line = line.strip()
            if line.startswith('#') or line.startswith('Slide'):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"title": line.lstrip('#').strip(), "points": [], "visual_suggestion": ""}
            elif line.startswith('-') or line.startswith('*'):
                if current_slide:
                    current_slide["points"].append(line.lstrip('-*').strip())
            elif line.startswith('Visual:') and current_slide:
                current_slide["visual_suggestion"] = line.replace('Visual:', '').strip()
        
        if current_slide:
            slides.append(current_slide)
        
        return slides
    
    except Exception as e:
        print(f"Error calling ChatGPT API: {str(e)}")
        # Return a basic structure if API call fails
        return [
            {"title": title, "points": ["Generated from text content"], "visual_suggestion": ""},
            {"title": "Main Points", "points": ["Please check the original text"], "visual_suggestion": ""}
        ]

def apply_presentation_style(prs, style):
    """Apply presentation style settings."""
    # Define style configurations
    style_configs = {
        "professional": {
            "bg_color": RGBColor(255, 255, 255),  # White
            "title_color": RGBColor(31, 73, 125),  # Dark blue
            "text_color": RGBColor(0, 0, 0),  # Black
            "accent_color": RGBColor(79, 129, 189)  # Light blue
        },
        "creative": {
            "bg_color": RGBColor(240, 240, 240),  # Light gray
            "title_color": RGBColor(192, 0, 0),  # Red
            "text_color": RGBColor(64, 64, 64),  # Dark gray
            "accent_color": RGBColor(255, 192, 0)  # Orange
        },
        "minimal": {
            "bg_color": RGBColor(255, 255, 255),  # White
            "title_color": RGBColor(0, 0, 0),  # Black
            "text_color": RGBColor(64, 64, 64),  # Dark gray
            "accent_color": RGBColor(192, 192, 192)  # Light gray
        }
    }
    
    # Use professional style as default if specified style is not found
    selected_style = style_configs.get(style.lower(), style_configs["professional"])
    
    # Apply style to slide master (this is limited in python-pptx)
    # For more comprehensive styling, consider using templates

def create_title_slide(prs, title, style):
    """Create the title slide."""
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Format title text
    title_paragraph = title_shape.text_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_run = title_paragraph.runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    
    # Add subtitle if there's a placeholder
    if slide.placeholders[1].has_text_frame:
        subtitle = slide.placeholders[1]
        subtitle.text = "Generated with AI"

def create_content_slide(prs, slide_data, style):
    """Create a content slide."""
    bullet_slide_layout = prs.slide_layouts[1]  # Bullet slide layout
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    # Set slide title
    title_shape = slide.shapes.title
    title_shape.text = slide_data["title"]
    
    # Format title
    title_paragraph = title_shape.text_frame.paragraphs[0]
    title_run = title_paragraph.runs[0]
    title_run.font.size = Pt(36)
    
    # Add bullet points
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()  # Clear default text
    
    # Add bullet points
    for point in slide_data["points"]:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0  # Top level bullet
        p.font.size = Pt(24)
    
    # Note about visual suggestions
    # In a full implementation, you might add placeholders or notes about the 
    # visual suggestions. Adding actual visuals would require more advanced handling.
    if slide_data.get("visual_suggestion"):
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = f"Visual suggestion: {slide_data['visual_suggestion']}"
